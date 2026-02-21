#!/usr/bin/env python3
"""
Simple document indexer for 1421 Foundation
No OpenAI, no embeddings - just SQLite with full-text search
"""

import os
import sys
import sqlite3
import hashlib
import re
from pathlib import Path
from datetime import datetime
import argparse

# Configure paths
BASE_DIR = Path(__file__).parent.parent.parent
DATA_DIR = BASE_DIR / "data"
DB_PATH = DATA_DIR / "knowledge_base.db"

# Supported file extensions
SUPPORTED_EXTS = {
    '.txt', '.md', '.csv', '.json', '.html', '.htm',
    '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt'
}

def read_file_content(file_path: Path) -> str:
    """Read content from different file types."""
    ext = file_path.suffix.lower()
    
    try:
        # Text files
        if ext in ['.txt', '.md', '.csv', '.json', '.html', '.htm']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        
        # PDF files
        elif ext == '.pdf':
            try:
                import pypdf
                reader = pypdf.PdfReader(file_path)
                text = ''
                for page in reader.pages[:10]:  # First 10 pages
                    text += page.extract_text() + '\n'
                return text[:10000]
            except ImportError:
                return f"[PDF file: {file_path.name}]"
        
        # Word documents
        elif ext in ['.docx', '.doc']:
            try:
                import docx
                doc = docx.Document(file_path)
                return '\n'.join([para.text for para in doc.paragraphs[:100]])
            except ImportError:
                return f"[Word document: {file_path.name}]"
        
        # Excel files
        elif ext in ['.xlsx', '.xls']:
            try:
                import pandas as pd
                df = pd.read_excel(file_path, nrows=50)
                return df.to_string()
            except ImportError:
                return f"[Excel file: {file_path.name}]"
        
        # PowerPoint
        elif ext in ['.pptx', '.ppt']:
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides[:20]:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return '\n'.join(text)
            except ImportError:
                return f"[PowerPoint: {file_path.name}]"
        
        else:
            return f"[{ext} file: {file_path.name}]"
            
    except Exception as e:
        return f"[Error reading {file_path.name}: {str(e)[:50]}]"

def extract_metadata(file_path: Path, content: str) -> dict:
    """Extract metadata from file."""
    metadata = {
        "source_file": str(file_path),
        "file_name": file_path.name,
        "file_size": file_path.stat().st_size,
        "file_type": file_path.suffix.lower(),
        "folder": str(file_path.parent.relative_to(DATA_DIR)) if file_path.parent != DATA_DIR else "root",
        "last_modified": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
        "title": file_path.stem,
    }
    
    # Try to extract year
    year_match = re.search(r'\b(1[4-9]\d{2}|20\d{2})\b', file_path.name + content[:1000])
    metadata["year"] = int(year_match.group(1)) if year_match else 0
    
    # Try to extract author
    author_match = re.search(r'by[_\s-]([A-Za-z\s]+)', file_path.name.lower() + content[:500].lower())
    metadata["author"] = author_match.group(1).strip().title() if author_match else "Unknown"
    
    return metadata

def index_documents(reset=False):
    """Main indexing function."""
    print("\n" + "="*60)
    print("📚 1421 FOUNDATION - SIMPLE DOCUMENT INDEXER")
    print("="*60)
    
    # Create data directory if it doesn't exist
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    
    # Reset if requested
    if reset and DB_PATH.exists():
        DB_PATH.unlink()
        print("🔄 Database reset")
    
    # Find all supported files
    all_files = []
    for ext in SUPPORTED_EXTS:
        all_files.extend(DATA_DIR.rglob(f"*{ext}"))
    
    print(f"\n🔍 Found {len(all_files)} files to index")
    
    if not all_files:
        print("❌ No files found to index")
        return
    
    # Connect to database
    conn = sqlite3.connect(str(DB_PATH))
    cursor = conn.cursor()
    
    # Create tables
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS documents (
            id TEXT PRIMARY KEY,
            title TEXT,
            author TEXT,
            year INTEGER,
            type TEXT,
            content_preview TEXT,
            source_file TEXT UNIQUE,
            file_size INTEGER,
            folder TEXT,
            last_modified TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    """)
    
    # Create indexes
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_year ON documents(year)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_author ON documents(author)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_type ON documents(type)")
    
    # Clear existing data
    cursor.execute("DELETE FROM documents")
    
    # Statistics
    stats = {"indexed": 0, "failed": 0, "by_type": {}}
    
    print("\n📄 Indexing files...")
    
    for i, file_path in enumerate(all_files, 1):
        ext = file_path.suffix.lower()
        stats["by_type"][ext] = stats["by_type"].get(ext, 0) + 1
        
        print(f"  [{i}/{len(all_files)}] {file_path.name[:60]}...", end=" ")
        
        try:
            # Read content
            content = read_file_content(file_path)
            
            # Extract metadata
            metadata = extract_metadata(file_path, content)
            
            # Create unique ID
            doc_id = hashlib.md5(f"{file_path}_{metadata['last_modified']}".encode()).hexdigest()[:16]
            
            # Content preview
            preview = content[:500] + "..." if len(content) > 500 else content
            
            # Insert into database
            cursor.execute("""
                INSERT OR REPLACE INTO documents 
                (id, title, author, year, type, content_preview, 
                 source_file, file_size, folder, last_modified)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """, (
                doc_id,
                metadata['title'],
                metadata['author'],
                metadata['year'],
                ext,
                preview,
                str(file_path),
                metadata['file_size'],
                metadata['folder'],
                metadata['last_modified']
            ))
            
            stats["indexed"] += 1
            print("✅")
            
        except Exception as e:
            print(f"❌ ({str(e)[:50]})")
            stats["failed"] += 1
    
    conn.commit()
    conn.close()
    
    # Print results
    print("\n" + "="*60)
    print("📊 INDEXING COMPLETE")
    print("="*60)
    print(f"✅ Files indexed: {stats['indexed']}")
    print(f"❌ Files failed: {stats['failed']}")
    print("\n📁 By file type:")
    for ext, count in stats["by_type"].items():
        print(f"  {ext}: {count}")
    print(f"\n💾 Database: {DB_PATH}")
    print(f"📦 Size: {DB_PATH.stat().st_size / 1024:.1f} KB" if DB_PATH.exists() else "")

def search_documents(query: str, limit: int = 10):
    """Simple search function."""
    if not DB_PATH.exists():
        print("❌ Database not found. Run indexing first.")
        return
    
    conn = sqlite3.connect(str(DB_PATH))
    cursor = conn.cursor()
    
    # Simple text search
    cursor.execute("""
        SELECT title, author, year, type, content_preview
        FROM documents
        WHERE title LIKE ? OR author LIKE ? OR content_preview LIKE ?
        ORDER BY year DESC
        LIMIT ?
    """, (f'%{query}%', f'%{query}%', f'%{query}%', limit))
    
    results = cursor.fetchall()
    conn.close()
    
    return results

def main():
    parser = argparse.ArgumentParser(description="Simple document indexer")
    parser.add_argument("--reset", action="store_true", help="Reset database")
    parser.add_argument("--search", type=str, help="Search for documents")
    args = parser.parse_args()
    
    if args.search:
        results = search_documents(args.search)
        if results:
            print(f"\n🔍 Search results for '{args.search}':")
            for i, (title, author, year, type_, preview) in enumerate(results, 1):
                print(f"\n{i}. {title}")
                print(f"   Author: {author}, Year: {year}, Type: {type_}")
                print(f"   {preview[:150]}...")
        else:
            print(f"\n❌ No results found for '{args.search}'")
        return
    
    index_documents(reset=args.reset)

if __name__ == "__main__":
    main()